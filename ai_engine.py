import os
import json
import random
from google import genai
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# --- CONFIGURATION ---
API_KEY = os.getenv("GOOGLE_API_KEY")
client = None

if not API_KEY:
    print("❌ CRITICAL ERROR: GOOGLE_API_KEY is missing from .env")
else:
    try:
        # 1. Initialize the Client
        client = genai.Client(api_key=API_KEY)
        print("✅ Connected to Google Gemini (Using model: gemini-2.5-flash)")
    except Exception as e:
        print(f"⚠️ Error initializing Gemini: {e}")

# --- HUAWEI MINDSPORE IMPORT ---
try:
    import mindspore
    from mindspore import Tensor
    import mindspore.ops as ops

    MINDSPORE_AVAILABLE = True
    print(f"✅ Huawei MindSpore v{mindspore.__version__} is active on CPU.")
except ImportError:
    MINDSPORE_AVAILABLE = False
    print("⚠️ MindSpore not found. Running in fallback mode.")

# ---DIGITAL (Standard Python Libs) ---
try:
    from docx import Document
    from pptx import Presentation
    from pypdf import PdfReader
except ImportError:
    print("⚠️ Document libraries not found. Please run: pip install pypdf python-docx python-pptx")


def extract_digital_text(filepath):
    """ Fast extraction for digital files (Word, PPT, selectable PDFs) """
    ext = os.path.splitext(filepath)[1].lower()

    try:
        if ext == '.pdf':
            reader = PdfReader(filepath)
            return "\n".join([page.extract_text() or "" for page in reader.pages])

        elif ext == '.docx':
            doc = Document(filepath)
            return "\n".join([p.text for p in doc.paragraphs])

        elif ext == '.pptx':
            prs = Presentation(filepath)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)

        else:
            return ""

    except Exception as e:
        print(f"Digital Extraction Error: {e}")
        return ""


# --- ENGINE 2: OPTICAL (Future Work / Placeholder) ---
def extract_optical_text(filepath):
    """ Placeholder for MindSpore OCR """
    # In future, connect to Huawei Cloud OCR here for images/scans
    return "[OCR] This is simulated text from an image/scan. (OCR marked as Future Work)"


# --- THE ROUTER (Connects main.py to the right engine) ---
def extract_text_from_file(filepath):
    """ Decides whether to use Digital extraction or Optical extraction based on file type """

    # 1. Try Digital Extraction first (PDF, DOCX, PPTX)
    text = extract_digital_text(filepath)

    # 2. If text is empty (e.g. scanned PDF or Image), try Optical (Future Work)
    if not text.strip():
        print(f"⚠️ No text found in {filepath}. Attempting Optical Extraction...")
        return extract_optical_text(filepath)

    return text


# --- RAG: RETRIEVAL LOGIC ---
def find_best_context(user_query, full_text):
    """
    Retrieval Logic.
    """
    import numpy as np
    # 1. Pre-processing (Python)
    paragraphs = [p for p in full_text.split('\n\n') if len(p) > 50]
    if not paragraphs: return full_text[:2000]
    # 2. Scoring (Python)
    query_words = set(user_query.lower().split())
    scores = []
    for p in paragraphs:
        score = sum(1 for w in query_words if w in p.lower())
        scores.append(score)
    # 3. MindSpore Decision Making
    if MINDSPORE_AVAILABLE:
        # Convert Score List to MindSpore Tensor
        score_tensor = Tensor(np.array(scores), mindspore.float32)
        # Perform Computation
        best_idx_tensor = ops.argmax(score_tensor)
        # Convert Result back to Python
        best_idx = int(best_idx_tensor.asnumpy())
        print(f"⚡ STRICT MODE: Retrieved Context via MindSpore (Index {best_idx})")
    else:
        # Only runs if 'import mindspore' failed at the very top of the file
        print("⚠️ MindSpore Library Missing! Falling back to standard CPU logic.")
        best_idx = np.argmax(scores)
    return paragraphs[best_idx][:2500]


def ask_bot(user_question, full_text_history=None):
    """
    The main entry point for the frontend.
    Handles 'No File' vs 'With File' logic automatically.
    """
    if not client:
        return "⚠️ Error: AI Engine is not connected."

    # CASE 1: GENERAL CHAT (No File Uploaded)
    if not full_text_history:
        print("ℹ️ No file loaded. Using General Tutor Mode.")
        prompt = (
            "You are Chokhmah, a helpful and encouraging AI tutor. "
            "The user has NOT uploaded any course notes yet.\n\n"
            f"USER QUESTION: {user_question}\n\n"
            "INSTRUCTIONS:\n"
            "1. **BE COMPREHENSIVE:** Answer the student's question in detail. Do not give short, one-line answers.\n"
            "2. **TEACHING STYLE:** Explain concepts clearly, using examples if necessary.\n"
            "3. **REMINDER:** Gently remind the user they can upload a PDF to get answers specific to their curriculum.\n"
            "4. **FORMATTING:** Use clean Markdown (Bold key terms, bullet points).\n"
        )
        try:
            response = client.models.generate_content(
                model='gemini-2.5-flash',
                contents=prompt
            )
            return response.text
        except Exception as e:
            return f"Error: {e}"

    # CASE 2: STRICT RAG (File IS Uploaded)
    else:
        # 1. Retrieve Context using MindSpore
        context = find_best_context(user_question, full_text_history)

        # 2. Generate Answer with Context
        return generate_answer(context, user_question)

# --- AI GENERATION ---
def generate_answer(context, question):
    if not client: return "⚠️ Error: Google Client not active."

    # UPDATED PROMPT FOR VERBOSE ANSWERS
    prompt = (
        "You are Chokhmah, an intelligent, encouraging study companion.\n"
        "You are given excerpts from the student's course material.\n\n"

        f"--- CONTEXT START ---\n{context}\n--- CONTEXT END ---\n\n"

        "USER QUESTION:\n"
        f"{question}\n\n"

        "INSTRUCTIONS:\n"
        "1. **BE COMPREHENSIVE:** Do not just give a one-line answer. Explain the concept fully using the provided context. Break it down so a student can understand.\n"
        "2. **STRICT GROUNDING:** Use ONLY the information in the context above. Do not make up outside facts.\n"
        "3. **FORMATTING:** Use **Bold** for key terms and lists for steps.\n"
        "4. **MATH:** If there are formulas, show them clearly using LaTeX ($$).\n"
    )

    try:
        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=prompt
        )
        return response.text

    except Exception as e:
        print(f"❌ Gemini Error: {e}")
        return "⚠️ Could not connect to Google AI."

# --- QUIZ GENERATION ---
def generate_quiz_question(full_text, difficulty="Medium", custom_topic=""):
    if not client: return None

    paragraphs = [p for p in full_text.split('\n\n') if len(p) > 100]
    if not paragraphs: return None

    # Smart Selection
    selected_text = ""
    if custom_topic:
        relevant = [p for p in paragraphs if custom_topic.lower() in p.lower()]
        selected_text = random.choice(relevant)[:2000] if relevant else random.choice(paragraphs)[:2000]
    else:
        selected_text = random.choice(paragraphs)[:2000]

    difficulty_instr = "Simple and direct." if difficulty == "Easy" else "Complex and tricky."
    topic_instr = f"Focus on: '{custom_topic}'." if custom_topic else ""

    prompt = (
        f"Generate 1 multiple-choice question based on:\n'{selected_text}'\n"
        f"Difficulty: {difficulty}. {topic_instr} {difficulty_instr}\n"
        f"Return ONLY valid JSON with keys: 'question', 'options', 'answer'."
    )

    try:
        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=prompt,
            config={'response_mime_type': 'application/json'}
        )

        raw = response.text.replace("```json", "").replace("```", "").strip()
        data = json.loads(raw)

        # Normalize keys
        return {k.lower(): v for k, v in data.items()}

    except Exception as e:
        print(f"❌ Quiz Error: {e}")
        return None


# --- SUMMARY GENERATION ---
def generate_summary(full_text, topic=""):
    if not client: return "AI Engine not connected."

    truncated_text = full_text[:4000]

    if topic:
        prompt = f"Summarize the following text focusing specifically on '{topic}':\n\n{truncated_text}"
    else:
        prompt = (f"Summarize the following study material into 3-6 key bullet points, but"
                  f"if more key point can be made use do that:\n\n{truncated_text}")

    try:
        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=prompt
        )
        return response.text
    except Exception as e:
        return f"Error generating summary: {e}"